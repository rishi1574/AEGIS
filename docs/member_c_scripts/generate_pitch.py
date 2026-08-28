import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_pitch_deck():
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project VANGUARD"
    subtitle.text = "Adversarial Evolution & Generative Intelligence Shield\nSelf-Healing Payment Security for the GenAI Era."

    # Slide 2: Problem
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Problem (GenAI Fraud Landscape 2026)"
    tf = body_shape.text_frame
    tf.text = "Speed: Agentic Hijacking (AI agents making autonomous purchases)"
    p = tf.add_paragraph()
    p.text = "Scale: Transaction Fuzzing (millions of micro-mutations to bypass thresholds)"
    p = tf.add_paragraph()
    p.text = "Sophistication: Deepfake ATOs & Synthetic ID Bust-Outs"
    p = tf.add_paragraph()
    p.text = "The Flaw: Current models suffer from Concept Drift. They degrade as fraud evolves."

    # Slide 3: Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The VANGUARD Solution (Red vs. Blue)"
    tf = body_shape.text_frame
    tf.text = "Concept: A continuous, closed-loop adversarial simulation."
    p = tf.add_paragraph()
    p.text = "Red Team Engine: Autonomous AI agents generating novel fraud vectors."
    p = tf.add_paragraph()
    p.text = "Blue Team Engine: Real-time detection pipeline learning from zero-day attacks."
    p = tf.add_paragraph()
    p.text = "Result: A self-healing defense system."

    # Slide 4: USPs
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key USPs"
    tf = body_shape.text_frame
    tf.text = "1. Autonomous Red Teaming (Zero-Day Discovery)"
    p = tf.add_paragraph()
    p.text = "2. Explainable AI (SHAP Waterfall for compliance)"
    p = tf.add_paragraph()
    p.text = "3. Federated Mule Detection (Cross-bank intelligence)"
    p = tf.add_paragraph()
    p.text = "4. Protection for Agentic Commerce (Know Your Agent - KYA)"

    # Save
    out_dir = "docs/assets"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Project_AEGIS_Pitch.pptx")
    prs.save(out_path)
    print(f"✅ Generated Pitch Deck at: {out_path}")

if __name__ == "__main__":
    create_pitch_deck()
